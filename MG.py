import g4f
import telebot
from pptx import Presentation
from pptx.util import Pt
bot = telebot.TeleBot('6786106072:AAGz8bE4mNyJscUDrvARjYiXOCAyX5FukD4')
@bot.message_handler(content_types=["text"]) 
def handle_text(message):
  text = message.text
  response = g4f.ChatCompletion.create(
    model="gpt-3.5-turbo",
    messages=[{"role": "user", "content": text}],
    stream=True    
  )
  ll2 = ""
  for item in response:
    ll2 += item
    
  root = Presentation()
  #создание первого слайда
  first_slide_layout = root.slide_layouts[1]
  slide = root.slides.add_slide(first_slide_layout)
  # создание заголовка слайда
  slide.shapes.title.text = 'Текст заголовка'
  # создание подзаголовка слайда
  text_frame = slide.placeholders[1].text_frame
  text_frame.text = ll2
  # формирование формата текста
  p = text_frame.paragraphs[0]
  p.font.size = Pt(14) 
  p.font.name = 'Times New Roman'
  root.save('Example.pptx')
  bot.send_message(message.chat.id, ll2)
    
bot.polling(none_stop=True, interval=0)

/////
/////
/////
/////
/////







