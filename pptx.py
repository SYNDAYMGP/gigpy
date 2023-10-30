фул количество символов в подзоголовке = 1850





from pptx import Presentation
from pptx.util import Pt
root = Presentation()
first_slide_layout = root.slide_layouts[1]
slide = root.slides.add_slide(first_slide_layout)
# заголовок слайда
slide.shapes.title.text = 'Текст заголовка'

# подзаголовок слайда
text_frame = slide.placeholders[1].text_frame
text_frame.text = 'Текст подзаголовка'

# формирование формата текста
p = text_frame.paragraphs[0]
p.font.size = Pt(14) 
p.font.name = 'Times New Roman'

# сохранение презентации 
root.save('Example.pptx')




,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
from pptx import Presentation
from pptx.util import Pt
root = Presentation()

#создание первого слайда
first_slide_layout = root.slide_layouts[1]
slide = root.slides.add_slide(first_slide_layout)
# создание заголовка слайда
slide.shapes.title.text = 'Текст заголовка'
# создание подзаголовка слайда
text_frame = slide.placeholders[1].text_frame
text_frame.text = 'Текст подзаголовка'
# формирование формата текста
p = text_frame.paragraphs[0]
p.font.size = Pt(14) 
p.font.name = 'Times New Roman'


#создание второго слайда
first_slide_layout2 = root.slide_layouts[1]
slide = root.slides.add_slide(first_slide_layout2)
# создание заголовка слайда
slide.shapes.title.text = 'Текст заголовка'
# создание подзаголовка слайда
text_frame = slide.placeholders[1].text_frame
text_frame.text = 'Текст подзаголовка'
# формирование формата текста
p = text_frame.paragraphs[0]
p.font.size = Pt(14) 
p.font.name = 'Times New Roman'


# сохранение презентации 
root.save('Example.pptx')
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
,,,,,,,,,,,,,,,
import g4f
import telebot
from pptx import Presentation
from pptx.util import Pt


bot = telebot.TeleBot('6786106072:AAGz8bE4mNyJscUDrvARjYiXOCAyX5FukD4')
@bot.message_handler(content_types=["text"])
def handle_text(message):
  with open('1.txt', 'w') as file:
    file.write('')
  text = message.text
  response = g4f.ChatCompletion.create(
    model="gpt-3.5-turbo",
    messages=[{"role": "user", "content": 'составь мне план (структуру) для моего доклада на тему (' + text + ') план должен состоять из глав в ответ ты должен записать только название глав без каких либо твоих пояснений и дополнений , не забудь пронумеровать их вот так 1, 2, 3, 4, 5'}], 
    stream=True
  )
  for msg in response:
    with open('1.txt', 'a', encoding='utf-8') as file:
      file.write(msg)
  with open('1.txt', 'r', encoding='utf-8') as file:
    text = file.read()
  start = text.find('1') + 1
  end = text.find('2')
  extracted_text = text[start:end].strip()
  ww1 = extracted_text
  
  if ww1:
    response1 = g4f.ChatCompletion.create(
    model="gpt-3.5-turbo",
    messages=[{"role": "user", "content": ww1}], 
    stream=True
   )
    text = next(response1)
  for msg in response1:
     root = Presentation()
     #создание первого слайда
     first_slide_layout = root.slide_layouts[1]
     slide = root.slides.add_slide(first_slide_layout)
     # создание заголовка слайда
     slide.shapes.title.text = 'Текст заголовка'
     # создание подзаголовка слайда
     text_frame = slide.placeholders[1].text_frame
     text_frame.text = 'hello'
     # формирование формата текста
     p = text_frame.paragraphs[0]
     p.font.size = Pt(14) 
     p.font.name = 'Times New Roman'
     root.save('Example.pptx')
  bot.send_document(message.chat.id, open('Example.pptx', 'rb'))
bot.polling(none_stop=True, interval=0)









#деление текта пополам 
def get_half_words(filename):
    with open(filename, 'r') as file:
        content = file.read()
        words = content.split()
        half_length = len(words) // 2
        half_words = " ".join(words[:half_length])
        half_words2 = " ".join(words[half_length:])
        return half_words, half_words2

file_name = '121.txt'
half_words, half_words2 = get_half_words(file_name)
print(half_words)
print(half_words2)



#считаем кол глав и записываем в переменную
with open("1.txt", encoding='utf-8') as file:
  content = file.read()
  lines = content.split("\n")
  chapter_count = 0
  for line in lines:
    if line:
      if line[0].isdigit():
        chapter_count += 1
        chapter_variable_name = "ww" + str(chapter_count)
        exec(chapter_variable_name + ' = line')
  print("Количество глав:", chapter_count)
  for i in range(1, chapter_count + 1):
    chapter_variable_name = "ww" + str(i)
    # Убираем .decode()
    print(chapter_variable_name + ":", eval(chapter_variable_name))
    
    

